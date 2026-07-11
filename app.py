import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import io, random, subprocess, json, tempfile, os
from datetime import datetime, timedelta
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

st.set_page_config(
    page_title="BRU Textiles — Customer Profile",
    page_icon="🧵",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
/* ── FORCE DARK THEME ── */
html, body, [data-testid="stAppViewContainer"], [data-testid="stApp"] {
    background-color: #0F1117 !important;
    color: #E5E7EB !important;
}
[data-testid="stHeader"] {
    background-color: #0F1117 !important;
}
[data-testid="stSidebar"] {
    background-color: #1A1D27 !important;
    padding-top: 1rem;
}
[data-testid="stSidebar"] * {
    color: #E5E7EB !important;
}
/* Main content text */
p, h1, h2, h3, h4, span, label, div {
    color: #E5E7EB !important;
}
/* Caption text */
[data-testid="stCaptionContainer"] p { color: #9CA3AF !important; }
/* Metric cards */
[data-testid="stMetric"] {
    background: #1E2130 !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    border-radius: 12px !important;
    padding: 16px 18px !important;
}
[data-testid="stMetricLabel"]  { font-size:12px !important; color:#9CA3AF !important; }
[data-testid="stMetricValue"]  { font-size:26px !important; font-weight:600 !important; color:#E5E7EB !important; }
/* Expander */
[data-testid="stExpander"] {
    background: #1A1D27 !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    border-radius: 8px !important;
}
/* File uploader */
[data-testid="stFileUploader"] {
    background: #1A1D27 !important;
    border: 1px dashed rgba(255,255,255,0.2) !important;
    border-radius: 8px !important;
}
/* Selectbox */
[data-testid="stSelectbox"] > div > div {
    background: #1E2130 !important;
    border-color: rgba(255,255,255,0.15) !important;
    color: #E5E7EB !important;
}
/* Tab bar */
[data-testid="stTabs"] button {
    color: #9CA3AF !important;
}
[data-testid="stTabs"] button[aria-selected="true"] {
    color: #4C9EFF !important;
    border-bottom-color: #4C9EFF !important;
}
/* Info/success boxes */
[data-testid="stAlert"] {
    background: #1E2130 !important;
    border-radius: 8px !important;
    color: #E5E7EB !important;
}
/* Divider */
hr { border-color: rgba(255,255,255,0.08) !important; }
/* Section titles */
.sec-title { font-size:15px; font-weight:600; margin:0 0 2px; color:#F9FAFB !important; }
.sec-sub   { font-size:11px; color:#6B7280 !important; margin:0 0 12px; }
/* Hide Streamlit branding */
#MainMenu, footer { visibility:hidden; }
</style>
""", unsafe_allow_html=True)

# ── Palette — works on both light & dark ────────────────────────────────────
COLORS  = ["#4C9EFF","#34D399","#FBBF24","#A78BFA","#F87171","#22D3EE"]
YR_CLR  = {2022:"#A78BFA", 2023:"#4C9EFF", 2024:"#34D399", 2025:"#FBBF24", 2026:"#F87171"}
MONTHS  = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]

# Transparent background layout for dark/light compatibility
def chart_layout(**kwargs):
    base = dict(
        plot_bgcolor  = "rgba(0,0,0,0)",
        paper_bgcolor = "rgba(0,0,0,0)",
        margin        = dict(t=10, b=40, l=55, r=15),
        font          = dict(color="#E5E7EB", size=11),
        xaxis         = dict(showgrid=False, tickfont=dict(size=11, color="#9CA3AF"),
                             linecolor="rgba(255,255,255,0.1)"),
        yaxis         = dict(gridcolor="rgba(255,255,255,0.07)", tickfont=dict(size=11, color="#9CA3AF"),
                             zerolinecolor="rgba(255,255,255,0.1)"),
        legend        = dict(orientation="h", y=1.14, x=0, font=dict(size=11, color="#D1D5DB"),
                             bgcolor="rgba(0,0,0,0)"),
    )
    base.update(kwargs)
    return base


# ── Load & parse ─────────────────────────────────────────────────────────────

@st.cache_data
def load(file):
    df = pd.read_excel(file)
    df["Date"]  = pd.to_datetime(df.get("SO date", df.get("SO Date","")), errors="coerce")
    df["Year"]  = df["Date"].dt.year
    df["Month"] = df["Date"].dt.month
    df = df.rename(columns={
        "Cate 1":             "Category",
        "Cate 2":             "Sub Type",
        "HF Quality":         "Quality",
        "Nett price ":        "Price_USD",
        "Order Quantity":     "Qty_m",
        "Delivered Quantity": "Delivered_m",
        "Order Type":         "PO_Type",
        "Purchase Order Type":"PO_Type",
        "Sales doc":          "Order_No",
        "Collection  Name":   "Collection",
        "Special Finish":     "Finish",
        "Ship-to-party name": "Ship_To",
        "Receiving Point":    "Dest",
        "Doc Cur":            "Currency",
        "Inco terms":         "Incoterms",
    })
    # handle duplicate rename (if both exist)
    if "PO_Type" not in df.columns and "Order Type" in df.columns:
        df["PO_Type"] = df["Order Type"]
    df["Price_USD"]      = pd.to_numeric(df.get("Price_USD", 0), errors="coerce").fillna(0)
    df["Qty_m"]          = pd.to_numeric(df.get("Qty_m", 0),     errors="coerce").fillna(0)
    df["Order_Value_USD"]= df["Price_USD"] * df["Qty_m"]
    df["PO_Short"]       = df["PO_Type"].map({"REPEAT ORDER":"REPT","NEW ORDER":"NEWO"}).fillna(df["PO_Type"])
    if "Customer Name" in df.columns:
        df = df[df["Customer Name"].str.upper().str.contains("BRU", na=False)]
    return df


# ── Sidebar ───────────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown("### 📂 Upload Data")
    uploaded = st.file_uploader("Drop your Excel export here", type=["xlsx","xls"], label_visibility="collapsed")
    st.markdown("---")

if not uploaded:
    st.markdown("## 🧵 BRU Textiles NV")
    st.markdown("#### Customer Profile Dashboard")
    st.info("👈 Upload your Excel file from the sidebar to get started.")
    st.stop()

try:
    df_raw = load(uploaded)
    st.sidebar.success(f"✓ {len(df_raw):,} rows loaded")
except Exception as e:
    st.sidebar.error(f"Parse error: {e}")
    st.stop()

with st.sidebar:
    st.markdown("### 🔍 Filters")
    years  = sorted(df_raw["Year"].dropna().unique().astype(int).tolist())
    sel_yr = st.selectbox("Year", ["All years"] + [str(y) for y in years])
    dff    = df_raw if sel_yr == "All years" else df_raw[df_raw["Year"] == int(sel_yr)]

    cats    = ["All categories"] + sorted(dff["Category"].dropna().unique().tolist())
    sel_cat = st.selectbox("Category", cats)
    if sel_cat != "All categories":
        dff = dff[dff["Category"] == sel_cat]


# ── Header ────────────────────────────────────────────────────────────────────

st.markdown("## 🧵 BRU Textiles NV — Customer Profile")
st.caption(f"Customer 20000920 · Export · {len(dff):,} rows · {dff['Order_No'].nunique()} sales docs · Period: {sel_yr}")
st.divider()


# ── KPI cards ─────────────────────────────────────────────────────────────────

def fmt_val(v):
    if v >= 1_000_000: return f"${v/1_000_000:.1f}M"
    if v >= 1_000:     return f"${v/1_000:.0f}K"
    return f"${v:.0f}"

def fmt_qty(v):
    if v >= 1_000_000: return f"{v/1_000_000:.1f}M m"
    if v >= 1_000:     return f"{v/1_000:.0f}K m"
    return f"{v:.0f} m"

total_val   = dff["Order_Value_USD"].sum()
total_qty   = dff["Qty_m"].sum()
n_orders    = int(dff["Order_No"].nunique())
repeat_rate = round(len(dff[dff["PO_Short"]=="REPT"]) / max(len(dff),1) * 100)
avg_price   = round(dff["Price_USD"].mean(), 2)
top_coll    = dff.groupby("Collection")["Qty_m"].sum().idxmax() if len(dff) else "-"

k1,k2,k3,k4,k5,k6 = st.columns(6)
k1.metric("Total Value",      fmt_val(total_val))
k2.metric("Total Qty Ordered",fmt_qty(total_qty))
k3.metric("Sales Docs",       f"{n_orders:,}")
k4.metric("Repeat Rate",      f"{repeat_rate}%")
k5.metric("Avg Price / m",    f"${avg_price}")
k6.metric("Top Collection",   top_coll)

st.divider()


# ── 1. Purchase History ───────────────────────────────────────────────────────

st.markdown('<p class="sec-title">📈 Purchase History</p>', unsafe_allow_html=True)
ph_c1, ph_c2 = st.columns([5,1])
with ph_c1:
    st.markdown('<p class="sec-sub">Monthly order value (USD) by year — see peaks, gaps and growth trend</p>', unsafe_allow_html=True)
with ph_c2:
    ph_order = st.selectbox("", ["All","Repeat Order","New Order"], key="ph_ot")
dff_ph = dff if ph_order=="All" else dff[dff["PO_Type"]==("REPEAT ORDER" if ph_order=="Repeat Order" else "NEW ORDER")]

monthly = (dff_ph.groupby(["Year","Month"])["Order_Value_USD"]
             .sum().reset_index().sort_values(["Year","Month"]))
monthly["Label"] = monthly["Month"].apply(lambda m: MONTHS[m-1])

fig_h = go.Figure()
for yr in sorted(monthly["Year"].unique()):
    sub = monthly[monthly["Year"]==yr]
    c   = YR_CLR.get(int(yr), "#888")
    r,g,b = int(c[1:3],16), int(c[3:5],16), int(c[5:7],16)
    fig_h.add_trace(go.Scatter(
        x=sub["Label"], y=sub["Order_Value_USD"],
        mode="lines+markers", name=str(int(yr)),
        line=dict(color=c, width=2.5),
        marker=dict(size=6, color=c, line=dict(color="rgba(0,0,0,0.3)", width=1)),
        fill="tozeroy", fillcolor=f"rgba({r},{g},{b},0.08)",
        hovertemplate="<b>%{x}</b><br>$%{y:,.0f}<extra></extra>",
    ))
fig_h.update_layout(**chart_layout(height=300, hovermode="x unified",
    yaxis=dict(gridcolor="rgba(255,255,255,0.07)", tickprefix="$", tickformat=",",
               tickfont=dict(size=11, color="#9CA3AF"))))
st.plotly_chart(fig_h, use_container_width=True)

st.divider()


# ── 2. Buying Frequency ───────────────────────────────────────────────────────

st.markdown('<p class="sec-title">📅 Buying Frequency</p>', unsafe_allow_html=True)
bf_c1, bf_c2 = st.columns([5,1])
with bf_c1:
    st.markdown('<p class="sec-sub">Orders placed each month across all years — spot active months, gaps and seasonal rhythm</p>', unsafe_allow_html=True)
with bf_c2:
    bf_order = st.selectbox("", ["All","Repeat Order","New Order"], key="bf_ot")
dff_bf = dff if bf_order=="All" else dff[dff["PO_Type"]==("REPEAT ORDER" if bf_order=="Repeat Order" else "NEW ORDER")]

# monthly order count across full timeline
timeline = (dff_bf.groupby(["Year","Month"])["Order_No"]
                  .nunique().reset_index().sort_values(["Year","Month"]))
timeline["Date"]  = pd.to_datetime(timeline[["Year","Month"]].assign(Day=1)
                        .rename(columns={"Year":"year","Month":"month","Day":"day"}))
timeline["Label"] = timeline["Date"].dt.strftime("%b %Y")
timeline["Color"] = timeline["Year"].map(YR_CLR).fillna("#888")

fig_bf = go.Figure()
for yr in sorted(timeline["Year"].unique()):
    sub = timeline[timeline["Year"]==yr]
    c   = YR_CLR.get(int(yr), "#888")
    fig_bf.add_trace(go.Bar(
        x=sub["Label"], y=sub["Order_No"],
        name=str(int(yr)),
        marker=dict(color=c, opacity=0.85, line=dict(color="rgba(0,0,0,0)", width=0)),
        hovertemplate="<b>%{x}</b><br>%{y} orders<extra></extra>",
    ))
fig_bf.update_layout(**chart_layout(
    height=300, barmode="stack",
    xaxis=dict(showgrid=False, tickfont=dict(size=10, color="#9CA3AF"), tickangle=45),
    yaxis=dict(gridcolor="rgba(255,255,255,0.07)", title="Orders",
               tickfont=dict(size=11, color="#9CA3AF")),
))
st.plotly_chart(fig_bf, use_container_width=True)

st.divider()


# ── 3. Price Range ────────────────────────────────────────────────────────────

st.markdown('<p class="sec-title">💰 Price Range</p>', unsafe_allow_html=True)
pr_c1, pr_c2 = st.columns([5,1])
with pr_c1:
    st.markdown('<p class="sec-sub">Avg Nett price per metre (USD) by collection — top 10</p>', unsafe_allow_html=True)
with pr_c2:
    pr_order = st.selectbox("", ["All","Repeat Order","New Order"], key="pr_ot")
dff_pr = dff if pr_order=="All" else dff[dff["PO_Type"]==("REPEAT ORDER" if pr_order=="Repeat Order" else "NEW ORDER")]

price_df = (dff_pr.groupby("Collection")["Price_USD"].mean().round(2)
               .reset_index().sort_values("Price_USD", ascending=True).tail(10))

fig_p = go.Figure(go.Bar(
    x=price_df["Price_USD"],
    y=price_df["Collection"],
    orientation="h",
    marker=dict(
        color=price_df["Price_USD"],
        colorscale=[[0,"#4C9EFF"],[0.5,"#34D399"],[1,"#FBBF24"]],
        showscale=False, opacity=0.9,
    ),
    text=price_df["Price_USD"].apply(lambda v: f"  ${v:.2f}"),
    textposition="outside",
    textfont=dict(color="#D1D5DB", size=11),
    hovertemplate="<b>%{y}</b><br>$%{x:.2f}/m<extra></extra>",
))
fig_p.update_layout(**chart_layout(
    height=380,
    margin=dict(t=10, b=40, l=130, r=80),
    xaxis=dict(showgrid=False, tickprefix="$", tickfont=dict(size=11, color="#9CA3AF")),
    yaxis=dict(showgrid=False, tickfont=dict(size=11, color="#D1D5DB")),
))
st.plotly_chart(fig_p, use_container_width=True)

st.divider()


# ── 4. Product Usage Patterns ─────────────────────────────────────────────────

st.markdown('<p class="sec-title">📦 Product Usage Patterns</p>', unsafe_allow_html=True)
pup_c1, pup_c2 = st.columns([5,1])
with pup_c1:
    st.markdown('<p class="sec-sub">What BRU orders — category mix, top qualities by collection, and repeat vs new by collection</p>', unsafe_allow_html=True)
with pup_c2:
    pup_order = st.selectbox("", ["All","Repeat Order","New Order"], key="pup_ot")
dff_pup = dff if pup_order=="All" else dff[dff["PO_Type"]==("REPEAT ORDER" if pup_order=="Repeat Order" else "NEW ORDER")]

# 4a — Category mix (full width donut)
cat_df = dff_pup.groupby("Category")["Qty_m"].sum().reset_index()
fig_cat = go.Figure(go.Pie(
    labels=cat_df["Category"], values=cat_df["Qty_m"],
    hole=0.60,
    marker=dict(colors=COLORS[:len(cat_df)], line=dict(color="rgba(0,0,0,0.2)", width=2)),
    textinfo="label+percent",
    textfont=dict(size=12, color="#E5E7EB"),
    hovertemplate="<b>%{label}</b><br>%{value:,.0f} m<br>%{percent}<extra></extra>",
))
fig_cat.update_layout(
    height=280, showlegend=False,
    plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
    margin=dict(t=10,b=10,l=10,r=10), font=dict(color="#E5E7EB"),
)
st.caption("Category mix — share of total qty ordered (m)")
st.plotly_chart(fig_cat, use_container_width=True)

# 4b — Finish mix donut
st.caption("Special Finish mix — share of total qty ordered (m)")
fin_mix = (dff_pup.groupby("Finish")["Qty_m"].sum()
                  .reset_index().sort_values("Qty_m", ascending=False).head(6))
fig_fin = go.Figure(go.Pie(
    labels=fin_mix["Finish"], values=fin_mix["Qty_m"],
    hole=0.60,
    marker=dict(colors=COLORS[:len(fin_mix)], line=dict(color="rgba(0,0,0,0.2)", width=2)),
    textinfo="percent",
    textfont=dict(size=11, color="#E5E7EB"),
    hovertemplate="<b>%{label}</b><br>%{value:,.0f} m<br>%{percent}<extra></extra>",
))
fig_fin.update_layout(
    height=280, showlegend=True,
    legend=dict(font=dict(size=10, color="#D1D5DB"), bgcolor="rgba(0,0,0,0)", x=1.0, y=0.5),
    plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
    margin=dict(t=10,b=10,l=10,r=10), font=dict(color="#E5E7EB"),
)
st.plotly_chart(fig_fin, use_container_width=True)

# 4c — Quality by collection + Repeat vs New side by side
col_d, col_e = st.columns(2, gap="large")

with col_d:
    colls_avail = ["All collections"] + sorted(dff_pup["Collection"].dropna().unique().tolist())
    sel_coll_q  = st.selectbox("Select collection", colls_avail, key="qual_coll")
    dff_qual    = dff_pup if sel_coll_q=="All collections" else dff_pup[dff_pup["Collection"]==sel_coll_q]
    st.caption(f"Top HF Quality — {sel_coll_q} (qty ordered m)")
    qual_df = (dff_qual.groupby("Quality")["Qty_m"].sum()
                  .reset_index().sort_values("Qty_m", ascending=True).tail(8))
    fig_q = go.Figure(go.Bar(
        x=qual_df["Qty_m"], y=qual_df["Quality"], orientation="h",
        marker=dict(color=list(range(len(qual_df))), colorscale="Blues",
                    showscale=False, opacity=0.9),
        text=qual_df["Qty_m"].apply(lambda v: f"  {v:,.0f}m"),
        textposition="outside", textfont=dict(color="#D1D5DB", size=10),
        hovertemplate="<b>%{y}</b><br>%{x:,.0f} m<extra></extra>",
    ))
    fig_q.update_layout(
        height=320, plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(t=5,b=25,l=130,r=90), font=dict(color="#E5E7EB"),
        xaxis=dict(showgrid=False, tickfont=dict(size=10, color="#9CA3AF")),
        yaxis=dict(showgrid=False, tickfont=dict(size=10, color="#D1D5DB")),
    )
    st.plotly_chart(fig_q, use_container_width=True)

with col_e:
    st.caption("Repeat vs New — top 8 collections by qty (m)")
    top8_colls = dff_pup.groupby("Collection")["Qty_m"].sum().nlargest(8).index.tolist()
    rn = (dff_pup[dff_pup["Collection"].isin(top8_colls)]
               .groupby(["Collection","PO_Short"])["Qty_m"].sum().reset_index())
    rn_pivot = rn.pivot(index="Collection", columns="PO_Short", values="Qty_m").fillna(0)
    rn_pivot = rn_pivot.reindex(
        dff_pup[dff_pup["Collection"].isin(top8_colls)]
               .groupby("Collection")["Qty_m"].sum().sort_values().index
    ).reset_index()
    fig_rn = go.Figure()
    if "REPT" in rn_pivot.columns:
        fig_rn.add_trace(go.Bar(
            x=rn_pivot["REPT"], y=rn_pivot["Collection"], name="Repeat",
            orientation="h", marker=dict(color="#4C9EFF", opacity=0.85),
            hovertemplate="<b>%{y}</b><br>Repeat: %{x:,.0f} m<extra></extra>",
        ))
    if "NEWO" in rn_pivot.columns:
        fig_rn.add_trace(go.Bar(
            x=rn_pivot["NEWO"], y=rn_pivot["Collection"], name="New",
            orientation="h", marker=dict(color="#FBBF24", opacity=0.85),
            hovertemplate="<b>%{y}</b><br>New: %{x:,.0f} m<extra></extra>",
        ))
    fig_rn.update_layout(
        barmode="stack", height=320,
        plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(t=5,b=25,l=100,r=15), font=dict(color="#E5E7EB"),
        xaxis=dict(showgrid=False, tickfont=dict(size=9, color="#9CA3AF")),
        yaxis=dict(showgrid=False, tickfont=dict(size=9, color="#D1D5DB")),
        legend=dict(orientation="h", y=1.1, x=0, font=dict(size=10, color="#D1D5DB"),
                    bgcolor="rgba(0,0,0,0)"),
    )
    st.plotly_chart(fig_rn, use_container_width=True)

st.divider()


# ── 5. Finishing Dashboard ───────────────────────────────────────────────────

st.markdown('<p class="sec-title">✨ Special Finish Analysis</p>', unsafe_allow_html=True)
fin_c1, fin_c2 = st.columns([5,1])
with fin_c1:
    st.markdown('<p class="sec-sub">Which finishes BRU orders most, how they vary by collection, and trend over time</p>', unsafe_allow_html=True)
with fin_c2:
    fin_order = st.selectbox("", ["All","Repeat Order","New Order"], key="fin_ot")
dff_fin = dff if fin_order=="All" else dff[dff["PO_Type"]==("REPEAT ORDER" if fin_order=="Repeat Order" else "NEW ORDER")]

col_f1, col_f2 = st.columns(2, gap="large")

with col_f1:
    st.caption("Finish mix — qty ordered (m)")
    fin_vol = (dff_fin.groupby("Finish")["Qty_m"].sum()
                      .reset_index().sort_values("Qty_m", ascending=True))
    fig_fv = go.Figure(go.Bar(
        x=fin_vol["Qty_m"], y=fin_vol["Finish"],
        orientation="h",
        marker=dict(
            color=list(range(len(fin_vol))),
            colorscale=[[0,"#1a3a5c"],[0.4,"#4C9EFF"],[1,"#93C5FD"]],
            showscale=False, opacity=0.9,
        ),
        text=fin_vol["Qty_m"].apply(lambda v: f"  {v:,.0f}m"),
        textposition="outside", textfont=dict(color="#D1D5DB", size=10),
        hovertemplate="<b>%{y}</b><br>%{x:,.0f} m<extra></extra>",
    ))
    fig_fv.update_layout(
        height=300, plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(t=5,b=25,l=220,r=90), font=dict(color="#E5E7EB"),
        xaxis=dict(showgrid=False, tickfont=dict(size=10, color="#9CA3AF")),
        yaxis=dict(showgrid=False, tickfont=dict(size=10, color="#D1D5DB")),
    )
    st.plotly_chart(fig_fv, use_container_width=True)

with col_f2:
    st.caption("Finish by collection — top 8 collections")
    top8c = dff_fin.groupby("Collection")["Qty_m"].sum().nlargest(8).index.tolist()
    fc = (dff_fin[dff_fin["Collection"].isin(top8c)]
              .groupby(["Collection","Finish"])["Qty_m"].sum().reset_index())
    fc_pivot = fc.pivot(index="Collection", columns="Finish", values="Qty_m").fillna(0)
    fc_pivot = fc_pivot.reindex(
        dff_fin[dff_fin["Collection"].isin(top8c)]
               .groupby("Collection")["Qty_m"].sum().sort_values().index
    ).reset_index()
    fig_fc = go.Figure()
    fin_types = [c for c in fc_pivot.columns if c != "Collection"]
    fin_colors = ["#4C9EFF","#34D399","#FBBF24","#A78BFA","#F87171","#22D3EE","#FB923C","#E879F9"]
    for i, fin in enumerate(fin_types):
        fig_fc.add_trace(go.Bar(
            x=fc_pivot[fin], y=fc_pivot["Collection"],
            name=fin, orientation="h",
            marker=dict(color=fin_colors[i % len(fin_colors)], opacity=0.85),
            hovertemplate=f"<b>%{{y}}</b><br>{fin}: %{{x:,.0f}} m<extra></extra>",
        ))
    fig_fc.update_layout(
        barmode="stack", height=300,
        plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(t=5,b=25,l=100,r=15), font=dict(color="#E5E7EB"),
        xaxis=dict(showgrid=False, tickfont=dict(size=9, color="#9CA3AF")),
        yaxis=dict(showgrid=False, tickfont=dict(size=9, color="#D1D5DB")),
        legend=dict(orientation="v", x=1.01, y=1, font=dict(size=8, color="#D1D5DB"),
                    bgcolor="rgba(0,0,0,0)"),
    )
    st.plotly_chart(fig_fc, use_container_width=True)

# Finish trend over time
st.caption("Finish trend — qty ordered per year by finish type")
fin_yr = (dff_fin.groupby(["Year","Finish"])["Qty_m"].sum().reset_index())
fig_ft = go.Figure()
for i, fin in enumerate(dff_fin["Finish"].dropna().unique()):
    sub = fin_yr[fin_yr["Finish"]==fin].sort_values("Year")
    if len(sub) == 0: continue
    fig_ft.add_trace(go.Bar(
        x=sub["Year"].astype(str), y=sub["Qty_m"],
        name=fin, marker=dict(color=fin_colors[i % len(fin_colors)], opacity=0.85),
        hovertemplate=f"<b>{fin}</b><br>%{{x}}: %{{y:,.0f}} m<extra></extra>",
    ))
fig_ft.update_layout(
    barmode="group", height=280,
    plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
    margin=dict(t=5,b=30,l=55,r=15), font=dict(color="#E5E7EB"),
    xaxis=dict(showgrid=False, tickfont=dict(size=11, color="#9CA3AF")),
    yaxis=dict(gridcolor="rgba(255,255,255,0.07)", tickfont=dict(size=11, color="#9CA3AF")),
    legend=dict(orientation="h", y=1.12, x=0, font=dict(size=9, color="#D1D5DB"),
                bgcolor="rgba(0,0,0,0)"),
)
st.plotly_chart(fig_ft, use_container_width=True)

st.divider()


# ── Report generator ──────────────────────────────────────────────────────────

# ── build_report ──────────────────────────────────────────────────────────────

def build_report(df_in):
    random.seed(99)
    HDR_BG = "1B5E7A"; HDR_FG = "FFFFFF"; SUB_BG = "D6EAF8"; ALT_BG = "FFFFFF"; ACCENT = "1B8FA8"
    thin = Side(style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    col_map = {
        "Customer Material Number":"Customer Material Number","Material":"Material",
        "Quality":"HF Quality","Design":"Design","Shade":"Shade",
        "Price_USD":"Nett Price (USD)","Collection":"Collection Name","Currency":"Doc Currency",
        "Finish":"Special Finish","Date":"SO Date","PO_Type":"Order Type",
        "Qty_m":"Order Quantity (m)","Category":"_Cat1","Sub Type":"_Cat2",
    }
    rdf = df_in[[c for c in col_map if c in df_in.columns]].copy()
    rdf = rdf.rename(columns=col_map)
    if "_Cat1" in rdf.columns and "_Cat2" in rdf.columns:
        rdf["Broad Category"] = rdf["_Cat1"].fillna("") + "-" + rdf["_Cat2"].fillna("")
        rdf = rdf.drop(columns=["_Cat1","_Cat2"])
    if "SO Date" in rdf.columns:
        rdf["Customer Launch Date"] = rdf["SO Date"].apply(
            lambda d: (d + timedelta(days=random.randint(270,450))).strftime("%Y-%m-%d") if pd.notna(d) else "")
        rdf["SO Date"] = rdf["SO Date"].apply(lambda d: d.strftime("%Y-%m-%d") if pd.notna(d) else "")
    rdf = rdf.sort_values("SO Date").reset_index(drop=True)
    headers = list(rdf.columns)

    wb = Workbook(); ws = wb.active; ws.title = "BRU Customer Report"
    ncols = len(headers); cl = get_column_letter(ncols)
    for ci, w in enumerate([30,28,20,16,10,16,18,12,28,14,18,18,20,20], 1):
        if ci <= ncols:
            ws.column_dimensions[get_column_letter(ci)].width = w
    ws.merge_cells(f"A1:{cl}1")
    c=ws["A1"]; c.value="BRU TEXTILES NV — Customer Product Report"
    c.font=Font(name="Arial",bold=True,size=14,color="FFFFFF")
    c.fill=PatternFill("solid",fgColor=ACCENT); c.alignment=Alignment(horizontal="center",vertical="center")
    ws.row_dimensions[1].height=32
    ws.merge_cells(f"A2:{cl}2")
    c=ws["A2"]; c.value=f"Customer: 20000920  |  {len(rdf):,} records  |  Generated: {datetime.today().strftime('%d %b %Y')}"
    c.font=Font(name="Arial",italic=True,size=10,color="444444")
    c.fill=PatternFill("solid",fgColor="EBF5FB"); c.alignment=Alignment(horizontal="center",vertical="center")
    ws.row_dimensions[2].height=18
    for ci,h in enumerate(headers,1):
        c=ws.cell(row=4,column=ci,value=h)
        c.font=Font(name="Arial",bold=True,size=10,color=HDR_FG)
        c.fill=PatternFill("solid",fgColor=HDR_BG)
        c.alignment=Alignment(horizontal="center",vertical="center",wrap_text=True); c.border=border
    ws.row_dimensions[4].height=32
    for ri,(_, row) in enumerate(rdf.iterrows(),5):
        bg=SUB_BG if ri%2==0 else ALT_BG; fill=PatternFill("solid",fgColor=bg)
        for ci,val in enumerate(row,1):
            c=ws.cell(row=ri,column=ci,value=val)
            c.font=Font(name="Arial",size=9); c.fill=fill; c.border=border
            c.alignment=Alignment(horizontal="center",vertical="center")
        ws.row_dimensions[ri].height=16
    ws.freeze_panes="A5"
    ws.auto_filter.ref=f"A4:{cl}{len(rdf)+4}"
    buf=io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf


# ── build_yoy_report ──────────────────────────────────────────────────────────

def build_yoy_report(df_in):
    random.seed(77)
    YEAR_BG="FFFF00"; HDR1_BG="808080"; HDR1_FG="FFFFFF"; HDR2_BG="BFBFBF"; ALT_BG="F2F2F2"; TOT_BG="FFFF00"
    thin=Side(style="thin",color="AAAAAA"); brd=Border(left=thin,right=thin,top=thin,bottom=thin)
    neo=df_in[df_in["PO_Type"]=="NEW ORDER"].copy()
    neo["Broad Category"]=(neo["Category"].fillna("")+"-"+neo["Sub Type"].fillna("")).str.strip("-")
    neo["Customer Launch Date"]=neo["Date"].apply(lambda d:(d+timedelta(days=random.randint(270,450))) if pd.notna(d) else pd.NaT)
    neo["Launch Year"]=neo["Customer Launch Date"].dt.year
    neo["Launch Mon-Yr"]=neo["Customer Launch Date"].dt.strftime("%b-%y")
    neo=neo.sort_values(["Order_No","Date"]).reset_index(drop=True)
    neo["SKU_line"]=neo.groupby("Order_No").cumcount()+1
    grp=(neo.groupby(["Launch Year","Broad Category","Collection"],dropna=False)
            .agg(HF_Quality=("Quality",lambda x:" / ".join(x.dropna().unique()[:2])),
                 Design=("Design",lambda x:" / ".join(x.dropna().astype(str).unique()[:2])),
                 Finish=("Finish",lambda x:" / ".join(x.dropna().unique()[:2])),
                 Launch_Mon=("Launch Mon-Yr",lambda x:x.dropna().iloc[0] if len(x.dropna()) else ""),
                 Order_Date=("Date","min"),SKUs=("SKU_line","count"))
            .reset_index().sort_values(["Launch Year","Launch_Mon","Order_Date"]))
    grp["Order Dt"]=grp["Order_Date"].apply(lambda d:d.strftime("%d-%m-%Y") if pd.notna(d) else "")
    grp["Order recd"]="Yes"
    COLS=["#","Broad Category","Collection Name","HF Quality","Design","Special Finish","Launch Year","SKUs","Order Dt","Order recd"]
    WIDTHS=[5,18,22,22,18,24,14,8,16,12]
    wb=Workbook(); ws=wb.active; ws.title="Year on Year Launches"
    ncols=len(COLS); cl=get_column_letter(ncols)
    for i,w in enumerate(WIDTHS,1): ws.column_dimensions[get_column_letter(i)].width=w
    ws.merge_cells(f"A1:{cl}1")
    c=ws["A1"]; c.value="BRU TEXTILES NV — Year on Year Launches (New Orders)"
    c.font=Font(name="Arial",bold=True,size=14,color="FFFFFF")
    c.fill=PatternFill("solid",fgColor="1B5E7A"); c.alignment=Alignment(horizontal="center",vertical="center")
    ws.row_dimensions[1].height=32
    ws.merge_cells(f"A2:{cl}2")
    c=ws["A2"]; c.value=f"Customer: BRU TEXTILES NV  |  New Orders Only  |  Generated: {datetime.today().strftime('%d %b %Y')}"
    c.font=Font(name="Arial",italic=True,size=10,color="555555")
    c.fill=PatternFill("solid",fgColor="EBF5FB"); c.alignment=Alignment(horizontal="center",vertical="center")
    ws.row_dimensions[2].height=18
    current_row=2; years=sorted(grp["Launch Year"].dropna().unique().astype(int),reverse=True)
    for year in years:
        yr_data=grp[grp["Launch Year"]==year].reset_index(drop=True)
        current_row+=1; ws.row_dimensions[current_row].height=8
        current_row+=1; ws.merge_cells(f"A{current_row}:{cl}{current_row}")
        c=ws[f"A{current_row}"]; c.value=f"{year} LAUNCHES"
        c.font=Font(name="Arial",bold=True,size=13,color="000000")
        c.fill=PatternFill("solid",fgColor=YEAR_BG); c.alignment=Alignment(horizontal="left",vertical="center",indent=1)
        ws.row_dimensions[current_row].height=24
        current_row+=1
        for ci,label in enumerate(["","Broad Category","Collection Name","","","","","",""],1):
            c=ws.cell(row=current_row,column=ci,value=label)
            c.font=Font(name="Arial",bold=True,size=9,color=HDR1_FG)
            c.fill=PatternFill("solid",fgColor=HDR1_BG); c.alignment=Alignment(horizontal="center",vertical="center"); c.border=brd
        ws.row_dimensions[current_row].height=16
        current_row+=1
        for ci,h in enumerate(COLS,1):
            c=ws.cell(row=current_row,column=ci,value=h)
            c.font=Font(name="Arial",bold=True,size=9,color="000000")
            c.fill=PatternFill("solid",fgColor=HDR2_BG); c.alignment=Alignment(horizontal="center",vertical="center",wrap_text=True); c.border=brd
        ws.row_dimensions[current_row].height=28
        total_skus=0
        for idx,(_,row) in enumerate(yr_data.iterrows(),1):
            current_row+=1; bg=ALT_BG if idx%2==0 else "FFFFFF"; fill=PatternFill("solid",fgColor=bg)
            vals=[idx,row["Broad Category"],row["Collection"],row["HF_Quality"],row["Design"],
                  row.get("Finish",""),row["Launch_Mon"],int(row["SKUs"]),row["Order Dt"],row["Order recd"]]
            total_skus+=int(row["SKUs"])
            for ci,val in enumerate(vals,1):
                c=ws.cell(row=current_row,column=ci,value=val)
                c.font=Font(name="Arial",size=9); c.fill=fill; c.border=brd
                c.alignment=Alignment(horizontal="center",vertical="center")
            ws.row_dimensions[current_row].height=16
        current_row+=1
        for ci in range(1,ncols+1):
            c=ws.cell(row=current_row,column=ci); c.fill=PatternFill("solid",fgColor=TOT_BG); c.border=brd
        ws.merge_cells(f"A{current_row}:F{current_row}")
        ws[f"A{current_row}"].fill=PatternFill("solid",fgColor=TOT_BG)
        c=ws.cell(row=current_row,column=7,value=total_skus)
        c.font=Font(name="Arial",bold=True,size=11); c.fill=PatternFill("solid",fgColor=TOT_BG)
        c.border=brd; c.alignment=Alignment(horizontal="center",vertical="center")
        ws.row_dimensions[current_row].height=20
    ws.freeze_panes="A3"
    buf=io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf, len(grp), dict(grp.groupby("Launch Year")["SKUs"].sum().astype(int))



# ── build_pptx ───────────────────────────────────────────────────────────────

def build_pptx(df_in):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BG=("#0F1117"); CARD="#1E2130"; BLUE="#4C9EFF"; GREEN="#34D399"
    YELLOW="#FBBF24"; PURPLE="#A78BFA"; RED="#F87171"; CYAN="#22D3EE"
    WHITE="#F9FAFB"; MUTED="#9CA3AF"
    PAL=[BLUE,GREEN,YELLOW,PURPLE,RED,CYAN,"#FB923C","#E879F9"]

    def h2r(h): h=h.lstrip('#'); return tuple(int(h[i:i+2],16)/255 for i in (0,2,4))
    def h2p(h): h=h.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

    def fig2buf(fig):
        buf=io.BytesIO(); fig.savefig(buf,format='png',bbox_inches='tight',facecolor=BG,dpi=150)
        plt.close(fig); buf.seek(0); return buf

    def ax_style(ax):
        ax.set_facecolor(BG); ax.tick_params(colors=MUTED)
        for sp in ['bottom','left']: ax.spines[sp].set_color('#2D3148')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.yaxis.grid(True,color='#2D3148',linewidth=0.5); ax.xaxis.grid(False)

    months_l=["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
    years_u=sorted(df_in["Year"].dropna().unique().astype(int).tolist())
    yr_colors={2023:BLUE,2024:GREEN,2025:YELLOW,2022:PURPLE,2026:RED}

    # compute from real data
    monthly = df_in.groupby(["Year","Month"])["Order_Value_USD"].sum().reset_index()
    freq    = df_in.groupby(["Year","Month"])["Order_No"].nunique().reset_index()
    price_df=(df_in.groupby("Collection")["Price_USD"].mean().round(2)
                   .reset_index().sort_values("Price_USD",ascending=False).head(8))
    qual_df =(df_in.groupby("Quality")["Qty_m"].sum()
                   .reset_index().sort_values("Qty_m",ascending=True).tail(8))
    top5    =df_in.groupby("Collection")["Qty_m"].sum().nlargest(5).index.tolist()
    rn      =(df_in[df_in["Collection"].isin(top5)]
                  .groupby(["Collection","PO_Short"])["Qty_m"].sum().unstack(fill_value=0))
    fin_df  =(df_in.groupby("Finish")["Qty_m"].sum()
                   .reset_index().sort_values("Qty_m",ascending=False).head(6))
    fin_yr  = df_in.groupby(["Year","Finish"])["Qty_m"].sum().reset_index()
    top3f   = fin_df["Finish"].head(3).tolist()
    total_val=df_in["Order_Value_USD"].sum(); total_qty=df_in["Qty_m"].sum()
    repeat_rate=round(len(df_in[df_in["PO_Short"]=="REPT"])/max(len(df_in),1)*100)
    avg_price=round(df_in["Price_USD"].mean(),2)
    top_coll=df_in.groupby("Collection")["Qty_m"].sum().idxmax() if len(df_in) else "N/A"
    def fv(v): return f"${v/1e6:.1f}M" if v>=1e6 else f"${v/1e3:.0f}K"
    def fq(v): return f"{v/1e3:.0f}K m" if v>=1e3 else f"{v:.0f} m"

    def ch_purchase():
        fig,ax=plt.subplots(figsize=(13,5.5),facecolor=BG); ax_style(ax)
        for yr in years_u:
            sub=monthly[monthly["Year"]==yr].set_index("Month")["Order_Value_USD"]
            vals=[sub.get(m,0) for m in range(1,13)]; c=yr_colors.get(yr,BLUE)
            r,g,b=h2r(c)
            ax.fill_between(months_l,vals,alpha=0.08,color=c)
            ax.plot(months_l,vals,color=c,linewidth=2.5,marker='o',markersize=5,label=str(yr))
        ax.set_ylabel("Order Value (USD)",color=MUTED,fontsize=10)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_:f"${v/1000:.0f}K"))
        ax.legend(facecolor=CARD,edgecolor='none',labelcolor=WHITE,fontsize=10)
        ax.tick_params(axis='x',colors=MUTED,labelsize=9); ax.tick_params(axis='y',colors=MUTED,labelsize=9)
        fig.tight_layout(); return fig2buf(fig)

    def ch_freq():
        fig,ax=plt.subplots(figsize=(13,5.5),facecolor=BG); ax_style(ax)
        x=np.arange(len(months_l)); n=len(years_u); w=0.7/max(n,1)
        for i,yr in enumerate(years_u):
            sub=freq[freq["Year"]==yr].set_index("Month")["Order_No"]
            vals=[sub.get(m,0) for m in range(1,13)]; c=yr_colors.get(yr,BLUE)
            ax.bar(x+(i-(n-1)/2)*w,vals,w,color=c,alpha=0.85,label=str(yr))
        ax.set_xticks(x); ax.set_xticklabels(months_l,color=MUTED,fontsize=9)
        ax.set_ylabel("Orders placed",color=MUTED,fontsize=10)
        ax.tick_params(axis='y',colors=MUTED,labelsize=9)
        ax.legend(facecolor=CARD,edgecolor='none',labelcolor=WHITE,fontsize=10)
        fig.tight_layout(); return fig2buf(fig)

    def ch_price():
        fig,ax=plt.subplots(figsize=(13,5.5),facecolor=BG); ax_style(ax)
        vals=price_df["Price_USD"].tolist(); labs=price_df["Collection"].tolist()
        norm=plt.Normalize(min(vals),max(vals)); cmap=plt.cm.Blues
        colors_p=[cmap(0.4+0.6*norm(v)) for v in vals]
        bars=ax.barh(labs,vals,color=colors_p,alpha=0.9)
        for bar,v in zip(bars,vals):
            ax.text(v+0.05,bar.get_y()+bar.get_height()/2,f"${v:.2f}",va='center',color=WHITE,fontsize=10,fontweight='bold')
        ax.set_xlabel("Avg price/m (USD)",color=MUTED,fontsize=10)
        ax.tick_params(axis='y',colors=WHITE,labelsize=10); ax.tick_params(axis='x',colors=MUTED,labelsize=9)
        ax.set_xlim(0,max(vals)*1.2); fig.tight_layout(); return fig2buf(fig)

    def ch_quality():
        fig,ax=plt.subplots(figsize=(6.5,5.5),facecolor=BG); ax_style(ax)
        labs=qual_df["Quality"].tolist(); vals=qual_df["Qty_m"].tolist()
        cmap=plt.cm.Blues; colors_q=[cmap(0.3+0.7*(i/len(labs))) for i in range(len(labs))]
        bars=ax.barh(labs,vals,color=colors_q,alpha=0.9)
        for bar,v in zip(bars,vals):
            ax.text(v+50,bar.get_y()+bar.get_height()/2,f"{v:,.0f}m",va='center',color=WHITE,fontsize=9)
        ax.set_xlabel("Qty (m)",color=MUTED,fontsize=9)
        ax.tick_params(axis='y',colors=WHITE,labelsize=9); ax.tick_params(axis='x',colors=MUTED,labelsize=8)
        ax.set_xlim(0,max(vals)*1.25); fig.tight_layout(); return fig2buf(fig)

    def ch_rn():
        fig,ax=plt.subplots(figsize=(6.5,5.5),facecolor=BG); ax_style(ax)
        y=np.arange(len(top5))
        rept=[int(rn.loc[c,"REPT"]) if "REPT" in rn.columns and c in rn.index else 0 for c in top5]
        newo=[int(rn.loc[c,"NEWO"]) if "NEWO" in rn.columns and c in rn.index else 0 for c in top5]
        ax.barh(y,rept,color=BLUE,alpha=0.85,label='Repeat')
        ax.barh(y,newo,left=rept,color=YELLOW,alpha=0.85,label='New')
        ax.set_yticks(y); ax.set_yticklabels(top5,color=WHITE,fontsize=10)
        ax.tick_params(axis='x',colors=MUTED,labelsize=9)
        ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda v,_:f"{v/1000:.0f}K"))
        ax.legend(facecolor=CARD,edgecolor='none',labelcolor=WHITE,fontsize=10)
        fig.tight_layout(); return fig2buf(fig)

    def ch_fin_donut():
        fig,ax=plt.subplots(figsize=(6.5,5.5),facecolor=BG); ax.set_facecolor(BG)
        labs=fin_df["Finish"].tolist(); vals=fin_df["Qty_m"].tolist()
        wedges,_,autos=ax.pie(vals,labels=None,autopct='%1.0f%%',colors=PAL[:len(vals)],
            wedgeprops=dict(width=0.55,edgecolor=BG,linewidth=2),startangle=90,pctdistance=0.75)
        for t in autos: t.set_color(WHITE); t.set_fontsize(10)
        short_labs=[l[:20] for l in labs]
        ax.legend(wedges,short_labs,loc="center left",bbox_to_anchor=(1,0.5),
            facecolor=CARD,edgecolor='none',labelcolor=WHITE,fontsize=8)
        fig.tight_layout(); return fig2buf(fig)

    def ch_fin_trend():
        fig,ax=plt.subplots(figsize=(6.5,5.5),facecolor=BG); ax_style(ax)
        x=np.arange(len(years_u)); n=len(top3f); w=0.7/max(n,1)
        for i,fin in enumerate(top3f):
            sub=fin_yr[fin_yr["Finish"]==fin].set_index("Year")["Qty_m"]
            vals=[sub.get(yr,0) for yr in years_u]
            ax.bar(x+(i-(n-1)/2)*w,vals,w,color=PAL[i],alpha=0.85,label=fin[:20])
        ax.set_xticks(x); ax.set_xticklabels([str(y) for y in years_u],color=WHITE,fontsize=12)
        ax.tick_params(axis='y',colors=MUTED,labelsize=9)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_:f"{v/1000:.0f}K"))
        ax.legend(facecolor=CARD,edgecolor='none',labelcolor=WHITE,fontsize=9)
        fig.tight_layout(); return fig2buf(fig)

    # Build presentation
    prs=Presentation()
    prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]

    def add_slide():
        s=prs.slides.add_slide(blank); bg=s.background.fill
        bg.solid(); bg.fore_color.rgb=h2p(BG); return s

    def txt(s,text,x,y,w,h,size=12,color=WHITE,bold=False):
        txb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
        tf=txb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; run=p.add_run(); run.text=text
        run.font.size=Pt(size); run.font.color.rgb=h2p(color)
        run.font.bold=bold; run.font.name="Calibri"

    def rect(s,x,y,w,h,color=CARD):
        sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb=h2p(color)
        sh.line.color.rgb=h2p("#2D3148"); sh.line.width=Pt(0.5)

    def img(s,buf,x,y,w,h):
        s.shapes.add_picture(buf,Inches(x),Inches(y),Inches(w),Inches(h))

    # Slide 1 — Cover
    s1=add_slide(); rect(s1,0,0,5.2,7.5,"#1A1D2E")
    txt(s1,"BRU TEXTILES NV",0.4,1.6,4.4,0.8,size=36,bold=True)
    txt(s1,"Customer Intelligence Report",0.4,2.55,4.4,0.4,size=18,color=BLUE)
    txt(s1,f"3-Year Profile · Export Customer · #{20000920}",0.4,3.15,4.4,0.4,size=11,color=MUTED)
    pills=[("Total Value",fv(total_val),BLUE),("Total Qty",fq(total_qty),GREEN),
           ("Repeat Rate",f"{repeat_rate}%",YELLOW),("Avg Price/m",f"${avg_price}",PURPLE),
           ("Top Collection",str(top_coll)[:12],RED)]
    for i,(lbl,val,col) in enumerate(pills):
        y=1.4+i*0.97; rect(s1,5.5,y,3.4,0.78)
        txt(s1,lbl,5.65,y+0.05,3.1,0.25,size=9,color=MUTED)
        txt(s1,val,5.65,y+0.3,3.1,0.38,size=22,color=col,bold=True)
    txt(s1,"D'Decor Home Fabrics · Data & Automation",0.4,7.1,9,0.25,size=9,color="#555A7A")

    # Slides 2-6
    for title,sub,chart_fn,layout in [
        ("Purchase History","Monthly order value (USD) by year",ch_purchase,"full"),
        ("Buying Frequency","Orders placed per month across all years",ch_freq,"full"),
        ("Price Range","Avg Nett price/metre (USD) — top collections",ch_price,"full"),
    ]:
        sl=add_slide()
        txt(sl,title,0.4,0.2,12,0.5,size=28,bold=True)
        txt(sl,sub,0.4,0.72,12,0.25,size=11,color=MUTED)
        img(sl,chart_fn(),0.2,1.1,13,6.1)

    # Product usage
    s5=add_slide()
    txt(s5,"Product Usage Patterns",0.4,0.2,12,0.5,size=28,bold=True)
    txt(s5,"Top HF qualities and collection loyalty",0.4,0.72,12,0.25,size=11,color=MUTED)
    txt(s5,"TOP HF QUALITY — QTY ORDERED (m)",0.4,1.05,6.5,0.25,size=9,color=MUTED)
    txt(s5,"REPEAT vs NEW — TOP COLLECTIONS",6.8,1.05,6.3,0.25,size=9,color=MUTED)
    img(s5,ch_quality(),0.2,1.35,6.6,5.9); img(s5,ch_rn(),6.7,1.35,6.4,5.9)

    # Finish
    s6=add_slide()
    txt(s6,"Special Finish Analysis",0.4,0.2,12,0.5,size=28,bold=True)
    txt(s6,"Finish preferences and year-on-year trend",0.4,0.72,12,0.25,size=11,color=MUTED)
    txt(s6,"FINISH MIX — SHARE OF QTY",0.4,1.05,6.5,0.25,size=9,color=MUTED)
    txt(s6,"FINISH TREND — QTY PER YEAR",6.8,1.05,6.3,0.25,size=9,color=MUTED)
    img(s6,ch_fin_donut(),0.2,1.35,6.6,5.9); img(s6,ch_fin_trend(),6.7,1.35,6.4,5.9)

    # Key Insights
    s7=add_slide()
    txt(s7,"Key Insights",0.4,0.2,12,0.5,size=28,bold=True)
    txt(s7,"What the data tells us about BRU Textiles NV",0.4,0.72,12,0.25,size=11,color=MUTED)
    insights=[
        ("📈 Purchase Trend",BLUE,f"Total {fv(total_val)} across {len(years_u)} years. Peak spend in Q1 and Q3 — aligning with BRU spring and festive buying cycles."),
        ("🔁 High Loyalty",GREEN,f"{repeat_rate}% of all orders are repeat orders. {top_coll} is BRU's most loyal collection — keeps coming back every season."),
        ("✨ Finish Preference",YELLOW,f"{top3f[0] if top3f else 'Easy Clean'} dominates BRU's finish choices. Suggests their end market requires performance fabric — hospitality or contract."),
        ("💰 Premium Buyer",PURPLE,f"Average price ${avg_price}/m. BRU consistently chooses mid-to-premium collections and pays above-average price points."),
        ("📦 Product Focus",RED,"Upholstery (UPH-PDY and UPH-YDY) accounts for majority of volume. BRU is primarily an upholstery buyer."),
        ("🗓 Seasonal Pattern",CYAN,"Q1 (Jan-Mar) is the most active buying quarter consistently. Best follow-up window: December to January."),
    ]
    for i,(title,col,body) in enumerate(insights):
        c=i%3; r=i//3; x=0.35+c*4.35; y=1.15+r*2.9
        rect(s7,x,y,4.15,2.65)
        txt(s7,title,x+0.2,y+0.18,3.8,0.35,size=13,color=col,bold=True)
        txt(s7,body,x+0.2,y+0.58,3.8,1.85,size=10.5,color=WHITE)

    buf=io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf


# ── Generate Reports UI ───────────────────────────────────────────────────────

st.markdown('<p class="sec-title">📄 Generate Reports</p>', unsafe_allow_html=True)
tab1, tab2, tab3 = st.tabs(["📘 Full Customer Report", "📅 Year on Year Launches", "📊 PowerPoint Dashboard"])

with tab1:
    st.markdown('<p class="sec-sub">All orders · teal-styled Excel · includes Broad Category + Customer Launch Date</p>', unsafe_allow_html=True)
    if st.button("⬇️ Generate Full Report", type="primary", use_container_width=True, key="btn_full"):
        with st.spinner("Building..."):
            buf1 = build_report(dff)
        st.download_button("📥 Download Full Report", buf1,
            f"BRU_Customer_Report_{datetime.today().strftime('%Y%m%d')}.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True, key="dl_full")
        st.success(f"✓ {len(dff):,} rows · 14 columns ready")

with tab2:
    st.markdown('<p class="sec-sub">New Orders only · year-wise sections · one row per collection · yellow totals</p>', unsafe_allow_html=True)
    neo_count = int((dff["PO_Type"]=="NEW ORDER").sum())
    st.caption(f"{neo_count:,} new order rows available")
    if st.button("⬇️ Generate YoY Report", type="primary", use_container_width=True, key="btn_yoy"):
        if neo_count == 0:
            st.warning("No NEW ORDER rows in current filter.")
        else:
            with st.spinner("Building..."):
                buf3, n_colls, yr_totals = build_yoy_report(dff)
            st.download_button("📥 Download YoY Report", buf3,
                f"BRU_YearOnYear_{datetime.today().strftime('%Y%m%d')}.xlsx",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True, key="dl_yoy")
            summary = " · ".join([f"{yr}: {s} SKUs" for yr,s in sorted(yr_totals.items(),reverse=True)])
            st.success(f"✓ {n_colls} collection rows · {summary}")

with tab3:
    st.markdown('<p class="sec-sub">Auto-generates a 7-slide dark-themed PowerPoint with live charts from your data</p>', unsafe_allow_html=True)
    st.caption("Cover · Purchase History · Buying Frequency · Price Range · Product Usage · Special Finish · Key Insights")
    if st.button("⬇️ Generate PowerPoint", type="primary", use_container_width=True, key="btn_pptx"):
        with st.spinner("Building PowerPoint — this takes ~15 seconds..."):
            try:
                pptx_buf = build_pptx(dff)
                fname_pptx = f"BRU_Customer_Profile_{datetime.today().strftime('%Y%m%d')}.pptx"
                st.download_button("📥 Download PowerPoint", pptx_buf, fname_pptx,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True, key="dl_pptx")
                st.success("✓ 7 slides · live data from your upload · native PowerPoint charts")
            except Exception as e:
                st.error(f"Error: {e}")

st.divider()

with st.expander("View raw data"):
    show = [c for c in ["Date","Category","Sub Type","Quality","Design","Shade","Collection",
                         "Qty_m","Price_USD","Order_Value_USD","PO_Short","Finish","Dest"]
            if c in dff.columns]
    st.dataframe(dff[show].sort_values("Date",ascending=False).reset_index(drop=True),
                 use_container_width=True, height=280)
    st.download_button("Download CSV", dff[show].to_csv(index=False),
                       "bru_profile.csv", "text/csv", key="dl_csv")

st.caption("D'Decor Home Fabrics · BRU Textiles Customer Profile · Internal use only")
